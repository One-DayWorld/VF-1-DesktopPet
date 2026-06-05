"""Generate ZAKU DESKTOP AGENT product pitch deck — v2 UI/UX optimised."""

import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ── Design Tokens ──────────────────────────────────────────────────────────────
BG       = RGBColor(0x0D, 0x0F, 0x0B)
BG2      = RGBColor(0x08, 0x0A, 0x07)
PANEL    = RGBColor(0x16, 0x19, 0x10)
PANEL2   = RGBColor(0x1E, 0x22, 0x16)
GREEN    = RGBColor(0x4B, 0x53, 0x20)
GREEN_LT = RGBColor(0x6E, 0x7A, 0x30)
GREEN_HL = RGBColor(0x92, 0xA2, 0x42)
ORANGE   = RGBColor(0xFF, 0x55, 0x00)
YELLOW   = RGBColor(0xFF, 0xD7, 0x00)
WHITE    = RGBColor(0xE8, 0xE8, 0xE0)
GRAY     = RGBColor(0x90, 0x98, 0x80)
GRAY_DK  = RGBColor(0x50, 0x58, 0x40)
BORDER   = RGBColor(0x3A, 0x40, 0x20)
BORDER2  = RGBColor(0x26, 0x2C, 0x14)
RED_DIM  = RGBColor(0x8B, 0x10, 0x00)
TEAL     = RGBColor(0x00, 0xCC, 0x77)   # scan / data accent

# ── Type Scale ─────────────────────────────────────────────────────────────────
T_CAP   = Pt(9)
T_SMALL = Pt(11)
T_BODY  = Pt(13)
T_MED   = Pt(16)
T_HEAD  = Pt(28)
T_TITLE = Pt(36)
T_HERO  = Pt(110)

W, H = Inches(13.33), Inches(7.5)

prs = Presentation()
prs.slide_width  = W
prs.slide_height = H

blank = prs.slide_layouts[6]


# ── Core Helpers ───────────────────────────────────────────────────────────────

def add_rect(slide, l, t, w, h, fill=None, line_color=None, line_w=Pt(1)):
    shape = slide.shapes.add_shape(1, l, t, w, h)
    shape.line.fill.background()
    if fill:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
    else:
        shape.fill.background()
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = line_w
    else:
        shape.line.fill.background()
    return shape


def add_text(slide, text, l, t, w, h,
             font_name="Courier New", size=None, bold=False, italic=False,
             color=WHITE, align=PP_ALIGN.LEFT, word_wrap=True):
    if size is None:
        size = T_BODY
    txb = slide.shapes.add_textbox(l, t, w, h)
    txb.word_wrap = word_wrap
    tf = txb.text_frame
    tf.word_wrap = word_wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = font_name
    run.font.size = size
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return txb


def add_para(tf, text, font_name="Courier New", size=None,
             bold=False, color=WHITE, align=PP_ALIGN.LEFT, space_before=Pt(10)):
    if size is None:
        size = T_BODY
    p = tf.add_paragraph()
    p.alignment = align
    p.space_before = space_before
    run = p.add_run()
    run.text = text
    run.font.name = font_name
    run.font.size = size
    run.font.bold = bold
    run.font.color.rgb = color
    return p


# ── Design Components ──────────────────────────────────────────────────────────

def bg(slide):
    add_rect(slide, 0, 0, W, H, fill=BG)


def corner_marks(slide):
    """HUD-style corner brackets with inner accent dot and coordinate labels."""
    s   = Inches(0.26)
    dot = Inches(0.07)
    pad = Inches(0.10)
    coords = [
        (pad, pad),
        (W - pad - s, pad),
        (pad, H - pad - s),
        (W - pad - s, H - pad - s),
    ]
    for lx, ly in coords:
        add_rect(slide, lx, ly, s, s, line_color=GREEN_LT, line_w=Pt(2))
        add_rect(slide, lx + (s - dot)/2, ly + (s - dot)/2, dot, dot, fill=GREEN_LT)
    # Coordinate labels in bottom corners (pure decoration)
    add_text(slide, "X:0000  Y:0000", Inches(0.10), H - Inches(0.24),
             Inches(2.0), Inches(0.20), size=Pt(7), color=GRAY_DK)
    add_text(slide, "X:1333  Y:0750", W - Inches(2.1), H - Inches(0.24),
             Inches(2.0), Inches(0.20), size=Pt(7), color=GRAY_DK, align=PP_ALIGN.RIGHT)


def scan_line(slide, y=None, color=GREEN):
    if y is None:
        y = Inches(1.02)
    add_rect(slide, Inches(0.35), y, W - Inches(0.7), Pt(1.5), fill=color)


def label_chip(slide, text, l, t, w=Inches(1.8), h=Inches(0.30),
               fill=GREEN, text_color=WHITE, size=None):
    if size is None:
        size = T_CAP
    add_rect(slide, l, t, w, h, fill=fill)
    add_text(slide, text, l + Inches(0.10), t + Pt(3), w, h,
             size=size, bold=True, color=text_color, align=PP_ALIGN.LEFT)


def section_header(slide, chip_text, title_text,
                   chip_color=GREEN, chip_text_color=WHITE, scan_color=None):
    add_rect(slide, 0, 0, W, Inches(1.02), fill=PANEL)
    add_rect(slide, 0, Inches(1.02), W, Pt(2), fill=scan_color or chip_color)
    label_chip(slide, chip_text, Inches(0.4), Inches(0.20),
               w=Inches(2.0), fill=chip_color, text_color=chip_text_color, size=T_CAP)
    add_text(slide, title_text, Inches(0.4), Inches(0.44), Inches(11), Inches(0.58),
             font_name="Arial Black", size=T_HEAD, bold=True, color=WHITE)


def bottom_bar(slide, label, page_n, total_n=9):
    add_rect(slide, 0, H - Inches(0.46), W, Inches(0.46),
             fill=RGBColor(0x10, 0x13, 0x0C))
    add_rect(slide, 0, H - Inches(0.46), W, Pt(1.5), fill=BORDER2)
    add_text(slide, f"ZAKU DESKTOP AGENT  //  {label}",
             Inches(0.4), H - Inches(0.37), W - Inches(2.0), Inches(0.32),
             size=Pt(8), color=GRAY_DK)
    # Page counter with accent box
    add_rect(slide, W - Inches(1.5), H - Inches(0.46),
             Inches(1.5), Inches(0.46), fill=PANEL)
    add_text(slide, f"{page_n:02d} / {total_n:02d}",
             W - Inches(1.4), H - Inches(0.38), Inches(1.3), Inches(0.32),
             size=Pt(9), bold=True, color=GREEN_LT, align=PP_ALIGN.RIGHT)


def feature_card(slide, l, t, w, h, number, title, body_lines,
                 title_color=ORANGE, accent=GREEN):
    add_rect(slide, l, t, w, h, fill=PANEL2, line_color=BORDER, line_w=Pt(1.2))
    # top + left accent bars
    add_rect(slide, l, t, w, Inches(0.07), fill=accent)
    add_rect(slide, l, t, Inches(0.05), h, fill=accent)
    # number badge
    badge_w, badge_h = Inches(0.38), Inches(0.26)
    add_rect(slide, l + Inches(0.12), t + Inches(0.10), badge_w, badge_h, fill=accent)
    add_text(slide, number, l + Inches(0.12), t + Inches(0.10), badge_w, badge_h,
             size=Pt(8), bold=True,
             color=BG if accent == YELLOW else WHITE, align=PP_ALIGN.CENTER)
    # title (supports \n for two-line titles)
    title_parts = title.split("\n") if "\n" in title else [title, ""]
    add_text(slide, title_parts[0], l + Inches(0.62), t + Inches(0.08),
             w - Inches(0.72), Inches(0.30), size=Pt(11), bold=True, color=title_color)
    if title_parts[1]:
        add_text(slide, title_parts[1], l + Inches(0.62), t + Inches(0.38),
                 w - Inches(0.72), Inches(0.22), size=Pt(10), bold=False, color=GRAY)
    # separator
    sep_y = t + Inches(0.65)
    add_rect(slide, l + Inches(0.12), sep_y, w - Inches(0.22), Pt(0.8), fill=BORDER)
    # body text
    txb = slide.shapes.add_textbox(l + Inches(0.14), sep_y + Inches(0.08),
                                   w - Inches(0.24), h - Inches(0.80))
    txb.word_wrap = True
    tf = txb.text_frame
    tf.word_wrap = True
    for i, line in enumerate(body_lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_before = Pt(3)
        run = p.add_run()
        run.text = line
        run.font.name = "Courier New"
        run.font.size = T_SMALL
        run.font.color.rgb = GRAY


def level_progress(slide, l, t, level, max_level=10, filled=GREEN, empty=BORDER2):
    """Mini level-progress bar strip."""
    seg_w = Inches(0.065)
    seg_h = Inches(0.10)
    gap   = Inches(0.018)
    for i in range(max_level):
        c = filled if i < level else empty
        add_rect(slide, l + i * (seg_w + gap), t, seg_w, seg_h, fill=c)


# ═════════════════════════════════════════════════════════════════════════════
# SLIDE 1 — Cover
# ═════════════════════════════════════════════════════════════════════════════
s1 = prs.slides.add_slide(blank)
bg(s1)

# Multi-density diagonal stripes (depth layers)
for i in range(24):
    x = Inches(-2.0) + i * Inches(0.78)
    a = 0x12 if i % 2 == 0 else 0x0F
    add_rect(s1, x, 0, Inches(0.28), H, fill=RGBColor(a, a + 3, a - 2))

# Faded unit watermark
add_text(s1, "MS-06F", Inches(6.2), Inches(2.2), Inches(7), Inches(3.2),
         font_name="Arial Black", size=Pt(120), bold=True,
         color=RGBColor(0x18, 0x1D, 0x11), align=PP_ALIGN.LEFT)

# Centre panel
add_rect(s1, Inches(0.9), Inches(0.95), Inches(11.5), Inches(5.45),
         fill=PANEL, line_color=GREEN, line_w=Pt(2.5))
add_rect(s1, Inches(0.9), Inches(0.95), Inches(11.5), Inches(0.10), fill=GREEN)
add_rect(s1, Inches(0.9), Inches(6.30), Inches(11.5), Inches(0.10), fill=ORANGE)

corner_marks(s1)

# Status chips
label_chip(s1, "MS-06F  ZAKU II", Inches(1.18), Inches(1.10),
           w=Inches(2.6), h=Inches(0.28), fill=GREEN, size=Pt(8))
label_chip(s1, "STATUS: ACTIVE", W - Inches(3.8), Inches(1.10),
           w=Inches(2.5), h=Inches(0.28), fill=TEAL,
           text_color=RGBColor(0x05, 0x08, 0x04), size=Pt(8))

# Main title
add_text(s1, "ZAKU", Inches(1.55), Inches(1.58), Inches(9), Inches(2.15),
         font_name="Arial Black", size=T_HERO, bold=True, color=WHITE, align=PP_ALIGN.LEFT)
add_text(s1, "DESKTOP AGENT", Inches(1.70), Inches(3.55), Inches(9.5), Inches(0.92),
         font_name="Arial Black", size=T_TITLE, bold=True, color=ORANGE, align=PP_ALIGN.LEFT)

# Tagline
add_text(s1, "你的 AI 桌面机甲伙伴  ·  ELECTRON  ·  THREE.JS  ·  AI MULTI-PROVIDER",
         Inches(1.70), Inches(4.62), Inches(9.5), Inches(0.46),
         size=Pt(12), color=GRAY, align=PP_ALIGN.LEFT)

# Version tag
add_text(s1, "v1.0  BUILD 2024", W - Inches(3.6), Inches(4.62),
         Inches(3.3), Inches(0.46), size=T_SMALL, color=GREEN_LT, align=PP_ALIGN.RIGHT)

# Targeting reticle (right-side HUD decoration)
add_rect(s1, Inches(11.0), Inches(2.6), Inches(0.06), Inches(1.5), fill=GREEN_LT)
add_rect(s1, Inches(10.3), Inches(3.32), Inches(1.5), Inches(0.06), fill=GREEN_LT)
add_rect(s1, Inches(11.5), Inches(2.4), Inches(0.55), Pt(1.5), fill=GREEN)
add_rect(s1, Inches(10.3), Inches(2.4), Inches(0.55), Pt(1.5), fill=GREEN)
add_text(s1, "TARGET LOCK", Inches(10.3), Inches(2.22), Inches(1.8), Inches(0.22),
         size=Pt(7), color=GREEN_LT, align=PP_ALIGN.CENTER)

bottom_bar(s1, "PRODUCT OVERVIEW  //  CONFIDENTIAL", 1)


# ═════════════════════════════════════════════════════════════════════════════
# SLIDE 2 — 项目概述
# ═════════════════════════════════════════════════════════════════════════════
s2 = prs.slides.add_slide(blank)
bg(s2)
corner_marks(s2)
section_header(s2, "OVERVIEW", "项目概述  PROJECT OVERVIEW")

# Left column — description
add_rect(s2, Inches(0.40), Inches(1.25), Inches(5.95), Inches(5.65),
         fill=PANEL2, line_color=BORDER, line_w=Pt(1))
add_rect(s2, Inches(0.40), Inches(1.25), Inches(5.95), Inches(0.07), fill=ORANGE)
add_rect(s2, Inches(0.40), Inches(1.25), Inches(0.06), Inches(5.65), fill=ORANGE)

add_text(s2, "WHAT IS ZAKU DESKTOP AGENT?",
         Inches(0.62), Inches(1.37), Inches(5.6), Inches(0.34),
         size=T_SMALL, bold=True, color=ORANGE)

txb = s2.shapes.add_textbox(Inches(0.62), Inches(1.82), Inches(5.55), Inches(4.80))
txb.word_wrap = True
tf = txb.text_frame
tf.word_wrap = True

desc_paras = [
    "Zaku Desktop Agent 是基于 Electron + Three.js\n的桌面 AI 伴侣应用。以经典 SD 扎古 II 机甲\n为形象，运行于 macOS 透明无边框悬浮窗口。",
    "机体常驻桌面角落，随时响应指令——\n对话、提醒、文件扫描、自动化工作流，\n一切尽在掌握。",
    "随着任务完成积累 XP，机体等级提升，\n肩盾、机炮、速射炮、推进器……\n武装配件逐级解锁，打造专属作战形态。",
]
for i, txt in enumerate(desc_paras):
    p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
    p.space_before = Pt(14) if i > 0 else Pt(0)
    run = p.add_run()
    run.text = txt
    run.font.name = "Courier New"
    run.font.size = T_BODY
    run.font.color.rgb = GRAY

# Right column — stats grid
stats = [
    ("PLATFORM",   "macOS · Electron 28",         GREEN_LT),
    ("RENDERER",   "Three.js r158 · WebGL",        GREEN_LT),
    ("AI ENGINE",  "OpenAI / Claude / Gemini",     TEAL),
    ("WINDOW",     "180 × 200 px 透明悬浮",        WHITE),
    ("XP SYSTEM",  "Lv.1 → Lv.10  十级成长",      ORANGE),
    ("WEAPONS",    "8 种武装配件逐级解锁",          ORANGE),
]
rx = Inches(6.55)
for i, (k, v, vc) in enumerate(stats):
    yy = Inches(1.25) + i * Inches(0.90)
    add_rect(s2, rx, yy, Inches(6.40), Inches(0.80), fill=PANEL2, line_color=BORDER)
    bar_c = GREEN if i < 2 else (TEAL if i == 2 else ORANGE)
    add_rect(s2, rx, yy, Inches(0.07), Inches(0.80), fill=bar_c)
    add_text(s2, k, rx + Inches(0.22), yy + Inches(0.06),
             Inches(2.0), Inches(0.26), size=T_CAP, bold=True, color=GRAY)
    add_text(s2, v, rx + Inches(0.22), yy + Inches(0.38),
             Inches(6.0), Inches(0.36), size=T_BODY, color=vc)

bottom_bar(s2, "PROJECT OVERVIEW", 2)


# ═════════════════════════════════════════════════════════════════════════════
# SLIDE 3 — 核心功能
# ═════════════════════════════════════════════════════════════════════════════
s3 = prs.slides.add_slide(blank)
bg(s3)
corner_marks(s3)
section_header(s3, "SYSTEMS", "核心功能模块  CORE SYSTEMS")

cards = [
    ("01", "AI COMBAT BRIDGE\n作战对话系统",
     ["· 支持 OpenAI / Claude / Gemini 多 Provider",
      "· 多轮上下文完整记忆",
      "· 机体心情随对话自然变化",
      "· 沉浸式 HUD 机甲风格对话界面"]),
    ("02", "MISSION CONTROL\n任务指挥系统",
     ["· 创建 / 编辑 / 完成提醒事项",
      "· 多任务列表分类管理",
      "· 完成任务获得 XP 积分奖励",
      "· 默认展示全部待完成任务"]),
    ("03", "FILE SCANNER\n文件扫描引擎",
     ["· 扫描桌面大体积冗余文件",
      "· AI 智能分析并给出删除建议",
      "· 一键 Reveal in Finder 定位",
      "· 安全确认后执行删除操作"]),
    ("04", "WORKFLOW ENGINE\n自动化引擎",
     ["· 自定义 Prompt 工作流模板",
      "· 一键触发 AI 批处理任务",
      "· 支持保存 / 编辑 / 删除",
      "· 结果实时流式回显"]),
    ("05", "PET CUSTOMIZATION\n机体定制系统",
     ["· 自定义机体名称与人格",
      "· 可更换头像 / 外观形象",
      "· 多情绪状态动画表现",
      "· 多 Provider API Key 管理"]),
    ("06", "GROWTH & LEVEL\n成长与升级",
     ["· XP 驱动十级成长体系",
      "· 情绪状态: 开心 / 兴奋 / 休眠",
      "· 跳跃动画庆祝任务完成",
      "· 等级与状态跨会话持久化"]),
]

cw, ch = Inches(4.15), Inches(2.58)
gx, gy = Inches(0.18), Inches(0.12)
mx, my = Inches(0.37), Inches(1.20)

for i, (num, title, body) in enumerate(cards):
    col = i % 3
    row = i // 3
    lx = mx + col * (cw + gx)
    ly = my + row * (ch + gy)
    feature_card(s3, lx, ly, cw, ch, num, title, body)

bottom_bar(s3, "CORE SYSTEMS", 3)


# ═════════════════════════════════════════════════════════════════════════════
# SLIDE 4 — AI 对话系统深潜
# ═════════════════════════════════════════════════════════════════════════════
s4 = prs.slides.add_slide(blank)
bg(s4)
corner_marks(s4)
section_header(s4, "AI SYSTEM", "AI 作战指挥桥  COMBAT BRIDGE",
               chip_color=ORANGE, chip_text_color=BG2, scan_color=ORANGE)

# Terminal mockup
add_rect(s4, Inches(0.40), Inches(1.22), Inches(5.80), Inches(5.70),
         fill=BG2, line_color=GREEN, line_w=Pt(2))
add_rect(s4, Inches(0.40), Inches(1.22), Inches(5.80), Inches(0.38), fill=GREEN)
add_text(s4, "  COMBAT BRIDGE — TERMINAL v2.6  [CONN: ACTIVE]",
         Inches(0.40), Inches(1.22), Inches(5.40), Inches(0.38),
         size=Pt(9), bold=True, color=BG2)
# Status dots (right side of terminal bar)
for i, c in enumerate([TEAL, YELLOW, GREEN_LT]):
    add_rect(s4, Inches(5.84) - Inches(0.30) - i * Inches(0.22),
             Inches(1.30), Inches(0.13), Inches(0.13), fill=c)

# Link established line
add_rect(s4, Inches(0.40), Inches(1.60), Inches(5.80), Pt(1), fill=BORDER2)
add_text(s4, "> ENCRYPTED LINK ESTABLISHED  //  UNIT MS-06F ONLINE",
         Inches(0.56), Inches(1.62), Inches(5.50), Inches(0.28),
         size=Pt(9), color=TEAL)

# Chat messages
chat_lines = [
    ("USER  >", WHITE,  "帮我整理今天的任务清单"),
    ("",        GRAY_DK, "─" * 40),
    ("ZAKU  >", ORANGE, "收到，指挥官。"),
    ("",        GRAY,   "已检测到 3 项待完成任务："),
    ("",        GRAY,   "  [!!] 提交项目报告  — 今日截止"),
    ("",        GRAY,   "  [ · ] 回复客户邮件 — 明日"),
    ("",        GRAY,   "  [ · ] 代码审查     — 本周"),
    ("ZAKU  >", ORANGE, "建议优先处理报告。"),
    ("",        ORANGE, "      需要我帮你起草摘要吗？"),
    ("USER  >", WHITE,  "好，帮我写"),
    ("ZAKU  >", TEAL,   "正在生成... ████████░░  80%"),
]
for i, (label, color, text) in enumerate(chat_lines):
    y = Inches(1.96) + i * Inches(0.36)
    if label:
        add_text(s4, label, Inches(0.56), y, Inches(0.95), Inches(0.34),
                 size=Pt(9), bold=True, color=color)
    add_text(s4, text, Inches(1.54), y, Inches(4.40), Inches(0.34),
             size=Pt(10), color=color)

# Right side feature cards
rx = Inches(6.42)
features = [
    (GREEN_LT, "01", "PROVIDER AGNOSTIC",
     "OpenAI · Anthropic Claude · Google Gemini\n随时热切换，无供应商锁定"),
    (ORANGE,   "02", "CONTEXT MEMORY",
     "多轮对话完整上下文记忆\n机体「记得」你上次说了什么"),
    (TEAL,     "03", "MOOD REACTION",
     "AI 响应后机体表达情绪\n跳跃动画 · 眼光扫描特效"),
    (YELLOW,   "04", "HUD INTERFACE",
     "机甲风 HUD 对话面板\n暗色主题 · 扫描线 · 绿光效果"),
]
for i, (color, num, title, body) in enumerate(features):
    yy = Inches(1.22) + i * Inches(1.40)
    add_rect(s4, rx, yy, Inches(6.52), Inches(1.26), fill=PANEL2, line_color=BORDER)
    add_rect(s4, rx, yy, Inches(0.08), Inches(1.26), fill=color)
    # number badge
    add_rect(s4, rx + Inches(0.18), yy + Inches(0.10),
             Inches(0.28), Inches(0.24), fill=color)
    add_text(s4, num, rx + Inches(0.18), yy + Inches(0.10),
             Inches(0.28), Inches(0.24), size=Pt(8), bold=True,
             color=BG2 if color == YELLOW else WHITE, align=PP_ALIGN.CENTER)
    add_text(s4, title, rx + Inches(0.58), yy + Inches(0.09),
             Inches(5.8), Inches(0.30), size=T_SMALL, bold=True, color=color)
    add_text(s4, body, rx + Inches(0.58), yy + Inches(0.46),
             Inches(5.8), Inches(0.74), size=T_BODY, color=GRAY)

bottom_bar(s4, "AI COMBAT BRIDGE", 4)


# ═════════════════════════════════════════════════════════════════════════════
# SLIDE 5 — 任务系统 & XP
# ═════════════════════════════════════════════════════════════════════════════
s5 = prs.slides.add_slide(blank)
bg(s5)
corner_marks(s5)
section_header(s5, "MISSION", "MISSION CONTROL — 任务指挥系统")

# Flow steps
flow = [
    ("01", "创建任务",  "输入任务文本\n设定截止日期\n分配到任务列表", GREEN),
    ("02", "执行追踪", "待完成视图优先\n过滤切换全部/完成\n实时状态同步",  ORANGE),
    ("03", "完成奖励", "标记完成获 XP\n机体跳跃庆祝\n状态持久化存储",    YELLOW),
]
for i, (num, title, body, color) in enumerate(flow):
    lx = Inches(0.45) + i * Inches(4.28)
    add_rect(s5, lx, Inches(1.30), Inches(4.0), Inches(2.60),
             fill=PANEL2, line_color=color, line_w=Pt(2))
    add_rect(s5, lx, Inches(1.30), Inches(4.0), Inches(0.10), fill=color)
    # number badge
    add_rect(s5, lx + Inches(0.18), Inches(1.46), Inches(0.50), Inches(0.48), fill=color)
    add_text(s5, num, lx + Inches(0.18), Inches(1.46), Inches(0.50), Inches(0.48),
             font_name="Arial Black", size=Pt(16), bold=True,
             color=BG2 if color == YELLOW else WHITE, align=PP_ALIGN.CENTER)
    add_text(s5, title, lx + Inches(0.86), Inches(1.47), Inches(3.0), Inches(0.48),
             font_name="Arial Black", size=T_MED, bold=True, color=color)
    add_text(s5, body, lx + Inches(0.22), Inches(2.06),
             Inches(3.60), Inches(1.65), size=T_BODY, color=GRAY)
    if i < 2:
        add_text(s5, "▶", lx + Inches(4.05), Inches(2.30),
                 Inches(0.30), Inches(0.40), size=Pt(18), color=BORDER, align=PP_ALIGN.CENTER)

# XP Progression table
add_rect(s5, Inches(0.45), Inches(4.12), Inches(12.42), Inches(2.75),
         fill=PANEL, line_color=BORDER)
add_rect(s5, Inches(0.45), Inches(4.12), Inches(12.42), Inches(0.09), fill=GREEN)
add_text(s5, "XP PROGRESSION MATRIX", Inches(0.65), Inches(4.26),
         Inches(5.5), Inches(0.30), size=Pt(10), bold=True, color=GREEN_LT)
add_text(s5, "— 完成任务积累 XP，解锁机体武装配件",
         Inches(5.4), Inches(4.26), Inches(7.2), Inches(0.30),
         size=Pt(9), color=GRAY)

xp_data = [
    ("Lv.1",  "0",    "",              WHITE,    None,    1),
    ("Lv.2",  "100",  "",              WHITE,    None,    2),
    ("Lv.3",  "300",  "尖刺肩盾",      GREEN_LT, GREEN,   3),
    ("Lv.4",  "600",  "120mm 肩炮",    GREEN_LT, GREEN,   4),
    ("Lv.5",  "1000", "MMP-80 速射炮", ORANGE,   ORANGE,  5),
    ("Lv.6",  "1500", "后置推进舱",    ORANGE,   ORANGE,  6),
    ("Lv.7",  "2200", "飞行翼",        ORANGE,   ORANGE,  7),
    ("Lv.8",  "3100", "指挥官天线",    YELLOW,   YELLOW,  8),
    ("Lv.9",  "4300", "全装甲胸板",    YELLOW,   YELLOW,  9),
    ("Lv.10", "6000", "夏亚定制限定",  RED_DIM,  RED_DIM, 10),
]
for i, (lv, xp, weapon, lc, bar_c, lvn) in enumerate(xp_data):
    cx = Inches(0.62) + i * Inches(1.22)
    bg_c = PANEL2 if weapon else PANEL
    add_rect(s5, cx, Inches(4.68), Inches(1.14), Inches(1.92), fill=bg_c)
    add_text(s5, lv, cx + Inches(0.04), Inches(4.72),
             Inches(1.06), Inches(0.28), size=Pt(10), bold=True,
             color=lc, align=PP_ALIGN.CENTER)
    add_text(s5, f"{xp} XP", cx + Inches(0.04), Inches(5.04),
             Inches(1.06), Inches(0.24), size=Pt(9), color=GRAY, align=PP_ALIGN.CENTER)
    if weapon:
        add_text(s5, weapon, cx, Inches(5.34), Inches(1.16), Inches(0.50),
                 size=Pt(9), color=lc, align=PP_ALIGN.CENTER)
    # progress bar at bottom of each cell
    if bar_c:
        add_rect(s5, cx + Inches(0.06), Inches(6.44), Inches(1.04), Inches(0.14), fill=bar_c)

bottom_bar(s5, "MISSION CONTROL", 5)


# ═════════════════════════════════════════════════════════════════════════════
# SLIDE 6 — 武器升级系统
# ═════════════════════════════════════════════════════════════════════════════
s6 = prs.slides.add_slide(blank)
bg(s6)
corner_marks(s6)
section_header(s6, "ARMAMENT", "武器升级系统  ARMAMENT UPGRADE",
               chip_color=ORANGE, chip_text_color=BG2, scan_color=ORANGE)

weapons = [
    (3,  "尖刺肩盾",   "Spiked Shoulder Shield",    "+X 右肩  ·  军绿装甲盾 + 六棱顶刺",        GREEN,   "300 XP"),
    (4,  "120mm 肩炮", "Shoulder Cannon",            "-X 左肩  ·  炮座 + 斜伸长管",               GREEN,   "600 XP"),
    (5,  "速射炮",     "MMP-80 Bazooka",             "+X 前臂  ·  斜持大口径炮管",                ORANGE,  "1000 XP"),
    (6,  "推进舱",     "Rear Thruster Pods ×2",      "背部双侧  ·  绿舱体 + 橙色发光喷嘴",        ORANGE,  "1500 XP"),
    (7,  "飞行翼",     "Flight Wing Binders ×2",     "两侧宽幅  ·  倾斜翼片构型",                 ORANGE,  "2200 XP"),
    (8,  "指挥天线",   "Commander V-fin",            "头顶黄色 V 形天线 + 侧感应器阵列",          YELLOW,  "3100 XP"),
    (9,  "全装甲",     "Full Armor Chest Plate",     "胸前厚重装甲板 + 双肩推进加力包",           YELLOW,  "4300 XP"),
    (10, "夏亚定制",   "Char's Custom Nozzle",       "底部红色高性能喷嘴 + 橙芯光效  [限定]",     RED_DIM, "6000 XP"),
]

for i, (lv, name_cn, name_en, desc, color, xp_req) in enumerate(weapons):
    col = i % 2
    row = i // 2
    lx = Inches(0.38) + col * Inches(6.54)
    ly = Inches(1.22) + row * Inches(1.44)
    ww, wh = Inches(6.22), Inches(1.32)

    add_rect(s6, lx, ly, ww, wh, fill=PANEL2, line_color=BORDER)
    add_rect(s6, lx, ly, Inches(0.07), wh, fill=color)

    # level badge
    add_rect(s6, lx + Inches(0.14), ly + Inches(0.10),
             Inches(0.68), Inches(0.28), fill=color)
    add_text(s6, f"Lv.{lv}", lx + Inches(0.14), ly + Inches(0.10),
             Inches(0.68), Inches(0.28), size=Pt(9), bold=True,
             color=BG2 if color == YELLOW else WHITE, align=PP_ALIGN.CENTER)

    # XP required (right-aligned badge)
    add_rect(s6, lx + ww - Inches(1.1), ly + Inches(0.10),
             Inches(0.98), Inches(0.24), fill=PANEL, line_color=BORDER)
    add_text(s6, xp_req, lx + ww - Inches(1.1), ly + Inches(0.10),
             Inches(0.98), Inches(0.24), size=Pt(8), color=GRAY, align=PP_ALIGN.CENTER)

    # Names
    add_text(s6, name_cn, lx + Inches(1.0), ly + Inches(0.07),
             Inches(2.1), Inches(0.36), font_name="Arial Black",
             size=Pt(14), bold=True, color=color)
    add_text(s6, name_en, lx + Inches(1.0), ly + Inches(0.46),
             Inches(2.7), Inches(0.26), size=Pt(10), color=GRAY)
    add_text(s6, desc, lx + Inches(1.0), ly + Inches(0.78),
             Inches(5.08), Inches(0.42), size=T_SMALL, color=WHITE)

    # Level progress bar (how close to max this weapon is)
    level_progress(s6, lx + Inches(0.14), ly + wh - Inches(0.16),
                   lv, filled=color)

bottom_bar(s6, "ARMAMENT UPGRADE SYSTEM", 6)


# ═════════════════════════════════════════════════════════════════════════════
# SLIDE 7 — 技术架构
# ═════════════════════════════════════════════════════════════════════════════
s7 = prs.slides.add_slide(blank)
bg(s7)
corner_marks(s7)
section_header(s7, "TECH STACK", "技术架构  TECHNICAL ARCHITECTURE")

layers = [
    ("PRESENTATION", "Electron BrowserWindow · 180×200 透明无边框 · macOS Always-on-top",
     GREEN_HL, Inches(1.14)),
    ("RENDER ENGINE", "Three.js r158 · GLTFLoader · SD Zaku II GLB · MeshStandardMaterial · AnimationMixer",
     GREEN_LT, Inches(1.96)),
    ("IPC BRIDGE",    "Electron contextBridge · preload.js · ipcRenderer ↔ ipcMain · 安全沙箱隔离",
     ORANGE,   Inches(2.78)),
    ("MAIN PROCESS",  "Node.js · electron-store 持久化 · AI Provider SDK · 文件系统操作",
     ORANGE,   Inches(3.60)),
    ("AI SERVICES",   "OpenAI API · Anthropic Claude API · Google Gemini API · Provider 热切换",
     TEAL,     Inches(4.42)),
    ("STORAGE",       "electron-store JSON · 宠物状态 · XP/等级 · 任务列表 · API Keys · 工作流",
     GRAY,     Inches(5.24)),
]

for idx, (title, body, color, y) in enumerate(layers):
    add_rect(s7, Inches(0.40), y, Inches(12.50), Inches(0.74),
             fill=PANEL2, line_color=BORDER)
    add_rect(s7, Inches(0.40), y, Inches(0.08), Inches(0.74), fill=color)
    # Layer index circle
    add_rect(s7, Inches(0.58), y + Inches(0.22), Inches(0.30), Inches(0.30),
             fill=PANEL, line_color=color, line_w=Pt(1))
    add_text(s7, str(idx + 1), Inches(0.58), y + Inches(0.22),
             Inches(0.30), Inches(0.30), size=Pt(9), bold=True,
             color=color, align=PP_ALIGN.CENTER)
    add_text(s7, title, Inches(1.05), y + Inches(0.08),
             Inches(3.0), Inches(0.28), size=Pt(10), bold=True, color=color)
    # Vertical separator
    add_rect(s7, Inches(4.06), y + Inches(0.20), Pt(1.5), Inches(0.34), fill=BORDER)
    add_text(s7, body, Inches(4.20), y + Inches(0.14),
             Inches(8.60), Inches(0.52), size=T_SMALL, color=WHITE)
    # Connector arrow to next layer
    if idx < len(layers) - 1:
        add_rect(s7, Inches(6.60), y + Inches(0.74), Pt(1), Inches(0.22), fill=BORDER)

bottom_bar(s7, "TECHNICAL ARCHITECTURE", 7)


# ═════════════════════════════════════════════════════════════════════════════
# SLIDE 8 — 机体情绪系统
# ═════════════════════════════════════════════════════════════════════════════
s8 = prs.slides.add_slide(blank)
bg(s8)
corner_marks(s8)
section_header(s8, "EMOTION AI", "机体情绪系统  EMOTION & ANIMATION",
               chip_color=ORANGE, chip_text_color=BG2, scan_color=ORANGE)

# No emoji — pure ASCII/text indicators
moods = [
    ("[ ◎ ]",  "HAPPY",    "开心",
     "默认待机状态\n轻微摇摆呼吸感\n眼部扫描光效",     GREEN,   "TRIGGER: IDLE"),
    ("[ !! ]", "EXCITED",  "兴奋",
     "完成任务触发\n持续跳跃动画\n4 秒后自动平复",       ORANGE,  "TRIGGER: TASK_COMPLETE"),
    ("[ Zzz ]","SLEEPING", "休眠",
     "长时间无操作\n缓慢呼吸下沉\n浮动休眠图标",         GRAY,    "TRIGGER: IDLE_TIMEOUT"),
    ("[ ◈ ]",  "BATTLE",   "战斗",
     "AI 对话时激活\n眼部红光加强\n机体前倾 3° 姿态",   RED_DIM, "TRIGGER: CHAT_OPEN"),
]

for i, (symbol, eng, chn, desc, color, trigger) in enumerate(moods):
    lx = Inches(0.40) + i * Inches(3.22)
    ly = Inches(1.25)
    ww, wh = Inches(3.02), Inches(5.68)

    add_rect(s8, lx, ly, ww, wh, fill=PANEL2, line_color=color, line_w=Pt(1.8))
    add_rect(s8, lx, ly, ww, Inches(0.08), fill=color)

    # Symbol display area
    add_rect(s8, lx + Inches(0.10), ly + Inches(0.20),
             ww - Inches(0.20), Inches(1.52), fill=BG2)
    # Status indicator dot
    add_rect(s8, lx + Inches(0.18), ly + Inches(0.28),
             Inches(0.12), Inches(0.12), fill=color)
    # ASCII symbol — large, centered
    add_text(s8, symbol, lx, ly + Inches(0.44), ww, Inches(0.85),
             size=Pt(30), bold=True, color=color, align=PP_ALIGN.CENTER)

    add_text(s8, eng, lx, ly + Inches(1.86), ww, Inches(0.40),
             font_name="Arial Black", size=Pt(15), bold=True,
             color=color, align=PP_ALIGN.CENTER)
    add_text(s8, chn, lx, ly + Inches(2.32), ww, Inches(0.36),
             size=T_BODY, color=WHITE, align=PP_ALIGN.CENTER)

    # Separator
    add_rect(s8, lx + Inches(0.20), ly + Inches(2.76),
             ww - Inches(0.40), Pt(1), fill=BORDER)

    add_text(s8, desc, lx + Inches(0.18), ly + Inches(2.90),
             ww - Inches(0.28), Inches(1.36),
             size=T_SMALL, color=GRAY, align=PP_ALIGN.CENTER)

    # Trigger label at bottom
    add_rect(s8, lx + Inches(0.10), ly + wh - Inches(0.44),
             ww - Inches(0.20), Inches(0.36), fill=PANEL)
    add_text(s8, trigger, lx + Inches(0.10), ly + wh - Inches(0.44),
             ww - Inches(0.20), Inches(0.36), size=Pt(8),
             color=GRAY, align=PP_ALIGN.CENTER)

bottom_bar(s8, "EMOTION & ANIMATION SYSTEM", 8)


# ═════════════════════════════════════════════════════════════════════════════
# SLIDE 9 — 结语 / CTA
# ═════════════════════════════════════════════════════════════════════════════
s9 = prs.slides.add_slide(blank)
bg(s9)

# Diagonal stripes
for i in range(24):
    x = Inches(-2.0) + i * Inches(0.78)
    a = 0x12 if i % 2 == 0 else 0x0F
    add_rect(s9, x, 0, Inches(0.28), H, fill=RGBColor(a, a + 3, a - 2))

# Faded watermark
add_text(s9, "DEPLOY", Inches(4.5), Inches(1.8), Inches(9), Inches(4),
         font_name="Arial Black", size=Pt(130), bold=True,
         color=RGBColor(0x18, 0x1C, 0x10), align=PP_ALIGN.LEFT)

# Centre panel
add_rect(s9, Inches(0.90), Inches(0.88), Inches(11.5), Inches(5.55),
         fill=PANEL, line_color=ORANGE, line_w=Pt(2.5))
add_rect(s9, Inches(0.90), Inches(0.88), Inches(11.5), Inches(0.10), fill=ORANGE)
add_rect(s9, Inches(0.90), Inches(6.33), Inches(11.5), Inches(0.10), fill=GREEN)

corner_marks(s9)

# Status chips
label_chip(s9, "DEPLOY ORDER", Inches(1.18), Inches(1.03),
           w=Inches(2.3), h=Inches(0.28), fill=ORANGE, text_color=BG2, size=Pt(8))
label_chip(s9, "STATUS: READY", W - Inches(3.7), Inches(1.03),
           w=Inches(2.5), h=Inches(0.28), fill=TEAL,
           text_color=RGBColor(0x04, 0x06, 0x03), size=Pt(8))

add_text(s9, "准备好部署了吗？",
         Inches(1.60), Inches(1.46), Inches(9.8), Inches(0.92),
         font_name="Arial Black", size=Pt(40), bold=True,
         color=WHITE, align=PP_ALIGN.LEFT)

add_text(s9, "Zaku Desktop Agent 现已可部署于任意 macOS 设备。",
         Inches(1.60), Inches(2.50), Inches(9.8), Inches(0.44),
         size=Pt(14), color=GRAY, align=PP_ALIGN.LEFT)

# Deploy command block
add_rect(s9, Inches(1.60), Inches(3.08), Inches(7.2), Inches(0.52),
         fill=BG2, line_color=TEAL, line_w=Pt(1))
add_text(s9, "> npm install  →  npm start  →  机体上线",
         Inches(1.80), Inches(3.10), Inches(7.0), Inches(0.48),
         size=Pt(13), color=TEAL, align=PP_ALIGN.LEFT)

# CTA buttons
btns = [
    ("▶  LAUNCH",    Inches(1.60),  GREEN,  BG2),
    ("◈  CONFIGURE", Inches(4.00),  ORANGE, BG2),
    ("★  UPGRADE",   Inches(6.50),  YELLOW, BG2),
]
for text, lx, fill, tc in btns:
    add_rect(s9, lx, Inches(3.82), Inches(2.15), Inches(0.58), fill=fill)
    add_text(s9, text, lx + Inches(0.15), Inches(3.88),
             Inches(1.92), Inches(0.48), font_name="Arial Black",
             size=Pt(13), bold=True, color=tc)

add_text(s9, "ELECTRON 28  ·  THREE.JS r158  ·  NODE.JS  ·  macOS  ·  AI MULTI-PROVIDER",
         Inches(1.60), Inches(4.65), Inches(9.8), Inches(0.38),
         size=Pt(10), color=BORDER, align=PP_ALIGN.LEFT)

add_text(s9, "MS-06F ZAKU II  —  FOR THE PRINCIPALITY OF ZEON  —  OPERATION: DESKTOP",
         Inches(1.60), Inches(5.10), Inches(9.8), Inches(0.36),
         size=Pt(10), color=GREEN_LT)

bottom_bar(s9, "DEPLOY READY", 9)


# ── Save ───────────────────────────────────────────────────────────────────────
out = os.path.join(os.path.expanduser("~"), "Documents", "DesktopPet", "ZAKU_DESKTOP_AGENT.pptx")
prs.save(out)
print(f"Saved: {out}")
